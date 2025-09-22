from pptx import Presentation

def parse_pptx(file_path: str) -> str:
    try:
        pres = Presentation(file_path)
        text = ""
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        # 암호화된 파일, 권한 문제, 손상된 파일 등 모든 예외 처리
        return ""